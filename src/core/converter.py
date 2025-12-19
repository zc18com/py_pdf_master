#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PDF转换模块 - 提供PDF与其他格式的转换功能
"""

import fitz
import io
import os
import tempfile
from typing import List, Optional, Dict, Any
from PIL import Image
from docx import Document
from docx.shared import Inches
import pdf2image
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
import pandas as pd
from pptx import Presentation
from pptx.util import Inches as PptInches


class PDFConverter:
    """PDF转换器类，提供PDF与其他格式的转换功能"""
    
    def __init__(self):
        self.pdf_parser = None
        self.dpi = 300  # 转换图片时的分辨率
    
    def set_parser(self, parser):
        """设置PDF解析器"""
        self.pdf_parser = parser
    
    def pdf_to_images(self, output_dir: str, format: str = "png", dpi: Optional[int] = None) -> List[str]:
        """
        PDF转换为图片
        
        Args:
            output_dir: 输出目录
            format: 图片格式 (png, jpeg, tiff)
            dpi: 图片分辨率
            
        Returns:
            List[str]: 生成的图片文件路径列表
        """
        try:
            if not self.pdf_parser or not self.pdf_parser.document:
                return []
            
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            # 设置DPI
            current_dpi = dpi if dpi else self.dpi
            
            # 获取PDF文件路径
            pdf_path = self.pdf_parser.file_path
            
            # 转换PDF为图片
            images = pdf2image.convert_from_path(
                pdf_path, 
                dpi=current_dpi,
                output_folder=output_dir,
                output_file="page",
                fmt=format,
                thread_count=4,
                use_pdftocairo=True
            )
            
            # 收集输出文件路径
            output_files = []
            for i in range(len(images)):
                output_file = os.path.join(output_dir, f"page_{i+1}.{format}")
                if os.path.exists(output_file):
                    output_files.append(output_file)
            
            return output_files
            
        except Exception as e:
            print(f"PDF转换为图片失败: {e}")
            return []
    
    def pdf_to_text(self, output_path: str) -> bool:
        """
        PDF转换为纯文本
        
        Args:
            output_path: 输出文本文件路径
            
        Returns:
            bool: 是否成功转换
        """
        try:
            if not self.pdf_parser or not self.pdf_parser.document:
                return False
            
            # 提取所有文本
            text = self.pdf_parser.extract_text()
            
            # 保存文本文件
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(text)
            
            return True
            
        except Exception as e:
            print(f"PDF转换为文本失败: {e}")
            return False
    
    def pdf_to_word(self, output_path: str) -> bool:
        """
        PDF转换为Word文档
        
        Args:
            output_path: 输出Word文件路径
            
        Returns:
            bool: 是否成功转换
        """
        try:
            if not self.pdf_parser or not self.pdf_parser.document:
                return False
            
            # 创建Word文档
            doc = Document()
            
            # 获取PDF文件
            doc_path = self.pdf_parser.file_path
            
            # 打开PDF
            doc_fitz = fitz.open(doc_path)
            
            for page_num in range(len(doc_fitz)):
                page = doc_fitz[page_num]
                
                # 提取文本
                text = page.get_text("text")
                
                # 如果有文本，添加到Word文档
                if text.strip():
                    doc.add_paragraph(text)
                
                # 提取图片
                images = page.get_images(full=True)
                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = doc_fitz.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # 创建临时图片文件
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_img:
                        temp_img.write(image_bytes)
                        temp_img_path = temp_img.name
                    
                    # 添加图片到Word文档
                    try:
                        doc.add_picture(temp_img_path, width=Inches(6))
                    finally:
                        # 删除临时文件
                        os.unlink(temp_img_path)
                
                # 每页添加分页符
                if page_num < len(doc_fitz) - 1:
                    doc.add_page_break()
            
            # 保存Word文档
            doc.save(output_path)
            doc_fitz.close()
            
            return True
            
        except Exception as e:
            print(f"PDF转换为Word失败: {e}")
            return False
    
    def pdf_to_excel(self, output_path: str) -> bool:
        """
        PDF转换为Excel文件
        
        Args:
            output_path: 输出Excel文件路径
            
        Returns:
            bool: 是否成功转换
        """
        try:
            if not self.pdf_parser or not self.pdf_parser.document:
                return False
            
            # 创建Excel工作簿
            excel_writer = pd.ExcelWriter(output_path, engine='openpyxl')
            
            # 获取PDF文件
            doc_path = self.pdf_parser.file_path
            
            # 打开PDF
            doc_fitz = fitz.open(doc_path)
            
            for page_num in range(len(doc_fitz)):
                page = doc_fitz[page_num]
                
                # 提取文本
                text = page.get_text("text")
                
                # 将文本分割成行
                lines = text.strip().split('\n')
                
                # 创建DataFrame
                df = pd.DataFrame(lines, columns=[f'第{page_num + 1}页'])
                
                # 将数据写入Excel工作表
                sheet_name = f'第{page_num + 1}页'
                df.to_excel(excel_writer, sheet_name=sheet_name, index=False)
            
            # 保存Excel文件
            excel_writer.close()
            doc_fitz.close()
            
            return True
            
        except Exception as e:
            print(f"PDF转换为Excel失败: {e}")
            return False
    
    def pdf_to_powerpoint(self, output_path: str) -> bool:
        """
        PDF转换为PowerPoint文件
        
        Args:
            output_path: 输出PowerPoint文件路径
            
        Returns:
            bool: 是否成功转换
        """
        try:
            if not self.pdf_parser or not self.pdf_parser.document:
                return False
            
            # 创建PowerPoint演示文稿
            prs = Presentation()
            
            # 获取PDF文件
            doc_path = self.pdf_parser.file_path
            
            # 打开PDF
            doc_fitz = fitz.open(doc_path)
            
            for page_num in range(len(doc_fitz)):
                page = doc_fitz[page_num]
                
                # 添加幻灯片
                slide_layout = prs.slide_layouts[5]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 提取文本
                text = page.get_text("text")
                
                # 如果有文本，添加到幻灯片
                if text.strip():
                    # 限制每页文本长度，避免PowerPoint格式问题
                    max_text_length = 500
                    if len(text) > max_text_length:
                        text = text[:max_text_length] + "..."
                    
                    # 添加文本框
                    left = PptInches(1)
                    top = PptInches(1)
                    width = PptInches(8)
                    height = PptInches(5)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = text
                
                # 提取图片
                images = page.get_images(full=True)
                for img_index, img in enumerate(images[:2]):  # 每页最多添加2张图片
                    xref = img[0]
                    base_image = doc_fitz.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # 创建临时图片文件
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_img:
                        temp_img.write(image_bytes)
                        temp_img_path = temp_img.name
                    
                    # 添加图片到幻灯片
                    try:
                        left = PptInches(1)
                        top = PptInches(2 + img_index * 2.5)
                        width = PptInches(8)
                        height = PptInches(2)
                        slide.shapes.add_picture(temp_img_path, left, top, width=width, height=height)
                    finally:
                        # 删除临时文件
                        os.unlink(temp_img_path)
            
            # 保存PowerPoint文件
            prs.save(output_path)
            doc_fitz.close()
            
            return True
            
        except Exception as e:
            print(f"PDF转换为PowerPoint失败: {e}")
            return False
    
    def images_to_pdf(self, image_paths: List[str], output_path: str, page_size: tuple = A4) -> bool:
        """
        多张图片转换为PDF
        
        Args:
            image_paths: 图片文件路径列表
            output_path: 输出PDF文件路径
            page_size: 页面大小 (width, height)
            
        Returns:
            bool: 是否成功转换
        """
        try:
            # 创建PDF画布
            c = canvas.Canvas(output_path, pagesize=page_size)
            
            for image_path in image_paths:
                if not os.path.exists(image_path):
                    continue
                
                # 打开图片
                img = Image.open(image_path)
                img_width, img_height = img.size
                
                # 计算图片在PDF页面上的缩放比例
                pdf_width, pdf_height = page_size
                scale = min(pdf_width / img_width, pdf_height / img_height)
                
                # 计算居中位置
                x = (pdf_width - img_width * scale) / 2
                y = (pdf_height - img_height * scale) / 2
                
                # 添加图片到PDF
                c.drawImage(
                    ImageReader(image_path),
                    x, y,
                    width=img_width * scale,
                    height=img_height * scale
                )
                
                # 结束当前页
                c.showPage()
            
            # 保存PDF
            c.save()
            
            return True
            
        except Exception as e:
            print(f"图片转换为PDF失败: {e}")
            return False
    
    def text_to_pdf(self, text_path: str, output_path: str, page_size: tuple = A4) -> bool:
        """
        文本文件转换为PDF
        
        Args:
            text_path: 输入文本文件路径
            output_path: 输出PDF文件路径
            page_size: 页面大小 (width, height)
            
        Returns:
            bool: 是否成功转换
        """
        try:
            # 读取文本内容
            with open(text_path, "r", encoding="utf-8") as f:
                text = f.read()
            
            # 创建PDF画布
            c = canvas.Canvas(output_path, pagesize=page_size)
            
            # 设置字体
            c.setFont("Helvetica", 12)
            
            # 设置页面参数
            pdf_width, pdf_height = page_size
            margin = 50
            line_height = 15
            max_lines_per_page = int((pdf_height - 2 * margin) / line_height)
            
            # 分割文本为行
            lines = []
            for paragraph in text.split("\n"):
                # 分割长行
                words = paragraph.split()
                current_line = ""
                for word in words:
                    test_line = current_line + (" " if current_line else "") + word
                    if len(test_line) > 80:  # 每行最大字符数
                        lines.append(current_line)
                        current_line = word
                    else:
                        current_line = test_line
                lines.append(current_line)
                lines.append("")  # 段落间空行
            
            # 添加文本到PDF
            current_page = 1
            current_line = 0
            
            for line in lines:
                if current_line >= max_lines_per_page:
                    # 新页面
                    c.showPage()
                    c.setFont("Helvetica", 12)
                    current_page += 1
                    current_line = 0
                
                y = pdf_height - margin - current_line * line_height
                c.drawString(margin, y, line)
                current_line += 1
            
            # 保存PDF
            c.save()
            
            return True
            
        except Exception as e:
            print(f"文本转换为PDF失败: {e}")
            return False
    
    def word_to_pdf(self, word_path: str, output_path: str) -> bool:
        """
        Word文档转换为PDF
        
        Args:
            word_path: 输入Word文件路径
            output_path: 输出PDF文件路径
            
        Returns:
            bool: 是否成功转换
        """
        try:
            # 使用docx2pdf库
            from docx2pdf import convert
            convert(word_path, output_path)
            return True
            
        except ImportError:
            print("docx2pdf库未安装")
            return False
        except Exception as e:
            print(f"Word转换为PDF失败: {e}")
            return False
    
    def pdf_to_html(self, output_path: str) -> bool:
        """
        PDF转换为HTML
        
        Args:
            output_path: 输出HTML文件路径
            
        Returns:
            bool: 是否成功转换
        """
        try:
            if not self.pdf_parser or not self.pdf_parser.document:
                return False
            
            # 获取PDF文件路径
            pdf_path = self.pdf_parser.file_path
            
            # 打开PDF
            doc = fitz.open(pdf_path)
            
            # 创建HTML内容
            html_content = "<html><head><title>PDF to HTML</title></head><body>\n"
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 提取HTML
                html = page.get_text("html")
                html_content += f"<div class='page'>{html}</div>\n"
            
            html_content += "</body></html>"
            
            # 保存HTML文件
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(html_content)
            
            doc.close()
            return True
            
        except Exception as e:
            print(f"PDF转换为HTML失败: {e}")
            return False
    
    def optimize_pdf(self, output_path: str, compress_images: bool = True, image_quality: int = 85) -> bool:
        """
        优化PDF文件（压缩）
        
        Args:
            output_path: 输出PDF文件路径
            compress_images: 是否压缩图片
            image_quality: 图片压缩质量 (0-100)
            
        Returns:
            bool: 是否成功优化
        """
        try:
            if not self.pdf_parser or not self.pdf_parser.document:
                return False
            
            # 获取PDF文件路径
            pdf_path = self.pdf_parser.file_path
            
            # 打开PDF
            doc = fitz.open(pdf_path)
            
            # 优化PDF
            if compress_images:
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    images = page.get_images(full=True)
                    
                    for img_index, img in enumerate(images):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # 压缩图片
                        try:
                            img = Image.open(io.BytesIO(image_bytes))
                            
                            # 转换为RGB
                            if img.mode != "RGB":
                                img = img.convert("RGB")
                            
                            # 保存压缩后的图片
                            img_byte_arr = io.BytesIO()
                            img.save(img_byte_arr, format="JPEG", quality=image_quality)
                            
                            # 替换原图片
                            doc.update_image(xref, img_byte_arr.getvalue())
                            
                        except Exception as e:
                            print(f"压缩图片失败: {e}")
                            continue
            
            # 保存优化后的PDF
            doc.save(output_path, garbage=4, deflate=True)
            doc.close()
            
            return True
            
        except Exception as e:
            print(f"优化PDF失败: {e}")
            return False


# 示例用法
if __name__ == "__main__":
    from src.core.pdf_parser import PDFParser
    
    # 创建PDF解析器和转换器
    parser = PDFParser()
    converter = PDFConverter()
    converter.set_parser(parser)
    
    # 打开PDF文件
    if parser.open_pdf("test.pdf"):
        # PDF转换为图片
        converter.pdf_to_images("output_images", format="png")
        
        # PDF转换为文本
        converter.pdf_to_text("output.txt")
        
        # 关闭PDF文件
        parser.close()
